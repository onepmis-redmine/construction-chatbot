import os
import json
import pandas as pd
from pathlib import Path
from datetime import datetime
from pdfminer.high_level import extract_text as pdfminer_extract_text
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
import chromadb
from sentence_transformers import SentenceTransformer
from tqdm import tqdm

# 디렉토리 경로
DOCS_DIR = Path("docs/")
RAW_TEXTS_DIR = Path("raw_texts/")
CHUNKS_DIR = Path("chunks/")
DB_DIR = Path("vector_db/")
EXCEL_PATH = Path("docs/qa_pairs.xlsx")
EXCEL_CHUNKS_PATH = Path("chunks/from_excel.json")
STATE_FILE = Path("pipeline_state.json")

# 디렉토리 생성
for directory in [RAW_TEXTS_DIR, CHUNKS_DIR, DB_DIR]:
    directory.mkdir(parents=True, exist_ok=True)

# 구성 요소 초기화
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=500,
    chunk_overlap=50,
    length_function=len
)
chroma_client = chromadb.PersistentClient(path=str(DB_DIR))
collection = chroma_client.get_or_create_collection(name="construction_manuals")
embedder = SentenceTransformer('all-MiniLM-L6-v2')

def load_state():
    """마지막 처리 상태(파일 수정 시간)를 로드합니다."""
    if STATE_FILE.exists():
        with open(STATE_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}

def save_state(state):
    """현재 상태를 저장합니다."""
    with open(STATE_FILE, "w", encoding="utf-8") as f:
        json.dump(state, f, ensure_ascii=False, indent=2)

def get_file_mtime(file_path):
    """파일 수정 시간을 문자열로 반환합니다."""
    return datetime.fromtimestamp(file_path.stat().st_mtime).isoformat()

def extract_docx_text(filepath):
    """DOCX 파일에서 텍스트를 추출합니다."""
    doc = Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_pptx_text(filepath):
    """PPTX 파일에서 텍스트를 추출합니다."""
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def extract_pdf_text(filepath):
    """PDF 파일에서 텍스트를 추출합니다."""
    return pdfminer_extract_text(filepath)

def extract_text_from_file(filepath):
    """파일 확장자에 따라 텍스트를 추출합니다."""
    filename = filepath.name
    if filename.endswith(".docx"):
        return extract_docx_text(filepath)
    elif filename.endswith(".pptx"):
        return extract_pptx_text(filepath)
    elif filename.endswith(".pdf"):
        return extract_pdf_text(filepath)
    else:
        print(f"⚠️ 지원하지 않는 파일 형식: {filename}")
        return None

def process_documents(state):
    """변경된 문서에서 텍스트를 추출해 raw_texts에 저장합니다."""
    updated = False
    for filepath in tqdm(list(DOCS_DIR.glob("*.[dp][od][cfx]*")), desc="텍스트 추출 중"):
        if filepath == EXCEL_PATH:
            continue  # 엑셀은 별도로 처리
        current_mtime = get_file_mtime(filepath)
        last_mtime = state.get(str(filepath), "")
        if current_mtime != last_mtime:
            text = extract_text_from_file(filepath)
            if text:
                output_path = RAW_TEXTS_DIR / (filepath.stem + ".txt")
                with open(output_path, "w", encoding="utf-8") as f:
                    f.write(text)
                state[str(filepath)] = current_mtime
                updated = True
    return updated

def chunk_texts(state):
    """원시 텍스트를 청크로 나누어 JSON으로 저장합니다."""
    updated = False
    for filepath in tqdm(list(RAW_TEXTS_DIR.glob("*.txt")), desc="텍스트 청킹 중"):
        current_mtime = get_file_mtime(filepath)
        last_mtime = state.get(str(filepath), "")
        output_path = CHUNKS_DIR / (filepath.stem + ".json")
        if current_mtime != last_mtime or not output_path.exists():
            with open(filepath, "r", encoding="utf-8") as f:
                text = f.read()
            chunks = text_splitter.split_text(text)
            output = [{"text": chunk} for chunk in chunks]
            with open(output_path, "w", encoding="utf-8") as f:
                json.dump(output, f, ensure_ascii=False, indent=2)
            state[str(filepath)] = current_mtime
            updated = True
    return updated

def process_excel(state):
    """엑셀 Q&A를 청크로 변환합니다."""
    if not EXCEL_PATH.exists():
        return False
    current_mtime = get_file_mtime(EXCEL_PATH)
    last_mtime = state.get(str(EXCEL_PATH), "")
    if current_mtime != last_mtime or not EXCEL_CHUNKS_PATH.exists():
        df = pd.read_excel(EXCEL_PATH)
        chunks = []
        for idx, row in df.iterrows():
            question = str(row["질문"]).strip()
            answer = str(row["답변"]).strip()
            source = str(row.get("출처", "excel_QA"))
            if question and answer:
                chunks.append({
                    "text": f"Q: {question}\nA: {answer}",
                    "source": source
                })
        with open(EXCEL_CHUNKS_PATH, "w", encoding="utf-8") as f:
            json.dump(chunks, f, ensure_ascii=False, indent=2)
        state[str(EXCEL_PATH)] = current_mtime
        return True
    return False

def embed_chunks(state):
    """청크를 벡터 데이터베이스에 임베딩합니다."""
    batch_size = 32
    total_added = 0
    updated = False
    for filepath in tqdm(list(CHUNKS_DIR.glob("*.json")), desc="청크 임베딩 중"):
        current_mtime = get_file_mtime(filepath)
        last_mtime = state.get(str(filepath), "")
        if current_mtime != last_mtime:
            try:
                with open(filepath, "r", encoding="utf-8") as f:
                    chunks = json.load(f)
            except json.JSONDecodeError:
                print(f"❌ 잘못된 JSON 형식: {filepath.name}")
                continue
            texts = [chunk.get("text", "").strip() for chunk in chunks if chunk.get("text")]
            if not texts:
                print(f"⚠️ 유효한 텍스트 없음: {filepath.name}")
                continue
            ids = [f"{filepath.stem}_{i}" for i in range(len(texts))]
            metadatas = [{"source": chunk.get("source", filepath.stem)} for chunk in chunks]
            embeddings = []
            for i in range(0, len(texts), batch_size):
                batch_texts = texts[i:i + batch_size]
                embeddings.extend(embedder.encode(batch_texts).tolist())
            if len(texts) != len(embeddings) or len(texts) != len(ids) or len(texts) != len(metadatas):
                print(f"❌ 데이터 길이 불일치: {filepath.name}")
                continue
            try:
                # 기존 항목 삭제로 중복 방지
                collection.delete(ids=ids)
                collection.add(
                    documents=texts,
                    embeddings=embeddings,
                    ids=ids,
                    metadatas=metadatas
                )
                total_added += len(texts)
                state[str(filepath)] = current_mtime
                updated = True
            except Exception as e:
                print(f"❌ 임베딩 오류: {filepath.name}: {str(e)}")
    print(f"✅ {total_added}개 문서 임베딩 완료")
    return updated

def main():
    """파이프라인을 실행합니다."""
    state = load_state()
    updated = False
    updated |= process_documents(state)
    updated |= chunk_texts(state)
    updated |= process_excel(state)
    updated |= embed_chunks(state)
    if updated:
        save_state(state)
        print("✅ 파이프라인 업데이트 완료")
    else:
        print("ℹ️ 업데이트 필요 없음")

if __name__ == "__main__":
    main()