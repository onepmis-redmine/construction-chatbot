import os
from pdfminer.high_level import extract_text as pdfminer_extract_text
from docx import Document
from pptx import Presentation
from tqdm import tqdm

# 폴더 경로
DOCS_DIR = "docs/"
RAW_TEXTS_DIR = "raw_texts/"

# 저장 폴더 없으면 만들기
os.makedirs(RAW_TEXTS_DIR, exist_ok=True)

# 텍스트 추출 함수들
def extract_docx_text(filepath):
    doc = Document(filepath)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_pptx_text(filepath):
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def my_extract_pdf_text(filepath):   # 함수 이름을 바꿨음
    return pdfminer_extract_text(filepath)

# 메인 추출 로직
for filename in tqdm(os.listdir(DOCS_DIR)):
    filepath = os.path.join(DOCS_DIR, filename)

    if filename.endswith(".docx"):
        text = extract_docx_text(filepath)
    elif filename.endswith(".pptx"):
        text = extract_pptx_text(filepath)
    elif filename.endswith(".pdf"):
        text = my_extract_pdf_text(filepath)  # 여기 호출할 때도 바꿔줬음
    else:
        print(f"지원하지 않는 파일 형식: {filename}")
        continue

    # 저장
    output_filename = os.path.splitext(filename)[0] + ".txt"
    output_path = os.path.join(RAW_TEXTS_DIR, output_filename)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
